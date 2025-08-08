import os
import logging
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches, Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import tempfile
import shutil
from PyPDF2 import PdfReader
import io

class PDFToPowerPointConverter:
    def __init__(self):
        self.temp_dir = None
    
    def convert(self, pdf_path, output_path, max_pages=30):
        """
        Convert PDF to PowerPoint presentation with editable text content
        
        Args:
            pdf_path (str): Path to the input PDF file
            output_path (str): Path for the output PowerPoint file
            max_pages (int): Maximum number of pages to process (default 30)
            
        Returns:
            bool: True if conversion successful, False otherwise
        """
        try:
            # Create temporary directory for images
            self.temp_dir = tempfile.mkdtemp()
            logging.info(f"Created temporary directory: {self.temp_dir}")
            
            # Check PDF page count first
            page_count = self.get_pdf_page_count(pdf_path)
            if page_count > max_pages:
                logging.warning(f"PDF has {page_count} pages, limiting to {max_pages} pages")
                last_page = max_pages
            else:
                last_page = page_count
            
            # Extract text from PDF
            pdf_texts = self.extract_text_from_pdf(pdf_path, last_page)
            
            # Convert PDF pages to images for visual reference
            logging.info(f"Converting PDF to images: {pdf_path} (pages 1-{last_page})")
            images = convert_from_path(
                pdf_path,
                dpi=150,  # Higher DPI for better quality when used as background
                output_folder=self.temp_dir,
                fmt='JPEG',
                thread_count=2,
                first_page=1,
                last_page=last_page
            )
            
            if not images:
                logging.error("No images generated from PDF")
                return False
            
            logging.info(f"Generated {len(images)} images from PDF")
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Use title and content layout for text-based slides
            title_layout = prs.slide_layouts[1]  # Title and Content layout
            blank_layout = prs.slide_layouts[6]  # Blank layout for image-only slides
            
            for i, image in enumerate(images):
                logging.info(f"Processing slide {i + 1}/{len(images)}")
                
                # Get extracted text for this page
                page_text = pdf_texts[i] if i < len(pdf_texts) else ""
                
                # Save image temporarily
                image_path = os.path.join(self.temp_dir, f"page_{i+1}.jpg")
                try:
                    if image.mode in ("RGBA", "P"):
                        image = image.convert("RGB")
                    image.save(image_path, 'JPEG', quality=90, optimize=True)
                except Exception as save_error:
                    logging.error(f"Error saving image {i+1}: {str(save_error)}")
                    image_path = os.path.join(self.temp_dir, f"page_{i+1}.png")
                    image.save(image_path, 'PNG')
                
                # Always create text-based slides with editable content
                slide = prs.slides.add_slide(blank_layout)
                
                # If we have extracted text, use it; otherwise create empty text boxes for user to edit
                if not page_text or len(page_text.strip()) < 20:
                    # Create placeholder text that user can replace
                    page_text = f"[Edit this text - Content from page {i+1}]\n\nClick here to add your content. You can replace this placeholder text with any content you want."
                
                self.create_editable_text_slide(slide, page_text, image_path, i+1)
            
            # Save PowerPoint presentation
            prs.save(output_path)
            logging.info(f"PowerPoint presentation saved: {output_path}")
            
            return True
            
        except Exception as e:
            logging.error(f"Error during PDF to PowerPoint conversion: {str(e)}")
            return False
        
        finally:
            # Clean up temporary directory
            if self.temp_dir and os.path.exists(self.temp_dir):
                try:
                    shutil.rmtree(self.temp_dir)
                    logging.info(f"Cleaned up temporary directory: {self.temp_dir}")
                except Exception as e:
                    logging.error(f"Error cleaning up temporary directory: {str(e)}")
    
    def extract_text_from_pdf(self, pdf_path, max_pages):
        """
        Extract text content from PDF pages
        
        Args:
            pdf_path (str): Path to the PDF file
            max_pages (int): Maximum number of pages to process
            
        Returns:
            list: List of text content for each page
        """
        texts = []
        try:
            with open(pdf_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                pages_to_process = min(len(pdf_reader.pages), max_pages)
                
                for i in range(pages_to_process):
                    try:
                        page = pdf_reader.pages[i]
                        text = page.extract_text()
                        # Clean up the text
                        text = text.replace('\n\n', '\n').strip()
                        texts.append(text)
                        logging.debug(f"Extracted {len(text)} characters from page {i+1}")
                    except Exception as e:
                        logging.warning(f"Error extracting text from page {i+1}: {str(e)}")
                        texts.append("")
                        
        except Exception as e:
            logging.error(f"Error reading PDF for text extraction: {str(e)}")
            
        return texts
    
    def create_editable_text_slide(self, slide, text, image_path, page_num):
        """
        Create a slide with fully editable text content - no images as main content
        """
        # Add title text box (always editable)
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = f"Page {page_num} - Click to Edit Title"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue-gray
        title_para.alignment = PP_ALIGN.CENTER
        
        # Add main content text box (large, editable area)
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(9), Inches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.margin_left = Inches(0.2)
        content_frame.margin_right = Inches(0.2)
        content_frame.margin_top = Inches(0.2)
        content_frame.margin_bottom = Inches(0.2)
        
        # Clean and format the text content
        if text.startswith("[Edit this text"):
            # This is placeholder text
            content_frame.text = text
            para = content_frame.paragraphs[0]
            para.font.size = Pt(16)
            para.font.color.rgb = RGBColor(127, 140, 141)  # Gray for placeholder
            para.font.italic = True
        else:
            # This is extracted text from PDF
            # Clean up the text and split into paragraphs
            cleaned_text = self.clean_extracted_text(text)
            paragraphs = cleaned_text.split('\n')
            
            for i, para_text in enumerate(paragraphs):
                if para_text.strip():
                    if i == 0:
                        p = content_frame.paragraphs[0]
                    else:
                        p = content_frame.add_paragraph()
                    
                    p.text = para_text.strip()
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.space_after = Pt(12)
                    p.line_spacing = 1.2
        
        # Add a small reference thumbnail at the bottom (optional, very small)
        try:
            with Image.open(image_path) as img:
                # Very small thumbnail for reference only - 15% of slide width
                thumb_width = Inches(1.5)
                thumb_height = Inches(1.1)
                
                # Position at bottom right corner
                left = Inches(8)
                top = Inches(6.2)
                
                slide.shapes.add_picture(image_path, left, top, thumb_width, thumb_height)
                
                # Add small text label for the thumbnail
                label_box = slide.shapes.add_textbox(
                    Inches(6.5), Inches(6.2),
                    Inches(1.4), Inches(0.3)
                )
                label_frame = label_box.text_frame
                label_frame.text = "Reference"
                label_para = label_frame.paragraphs[0]
                label_para.font.size = Pt(8)
                label_para.font.color.rgb = RGB